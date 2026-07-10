"""Builds a 5-7 slide portfolio health deck from the JSON reports main.py
already wrote to /outputs. Reads JSON only — no scoring logic here either;
this module's one piece of original content is the LLM-generated VP
recommendations slide.

Run:
    python presentation/generator.py
"""

import json
import re
import sys
from datetime import datetime
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from dotenv import load_dotenv

BASE_DIR = Path(__file__).parent.parent
OUTPUT_DIR = BASE_DIR / "outputs"
DECK_PATH = OUTPUT_DIR / "portfolio_health_deck.pptx"

load_dotenv(BASE_DIR / ".env")

sys.path.insert(0, str(BASE_DIR))
import llm_client
import rag_engine
from llm_explainer import plain_english_phrase

# Midnight Executive palette + RAG accents
NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x3A, 0x41, 0x5C)
MUTED = RGBColor(0x6B, 0x72, 0x8E)
CARD_BG = RGBColor(0xF3, 0xF5, 0xFB)

RAG_COLORS = {
    "Green": RGBColor(0x2E, 0x8B, 0x57),
    "Amber": RGBColor(0xE0, 0x9F, 0x2D),
    "Red": RGBColor(0xC0, 0x39, 0x2B),
    "Insufficient Data": RGBColor(0x80, 0x80, 0x80),
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _report_history() -> dict[str, list[dict]]:
    """Every stored weekly report per project, oldest to newest — the
    real, accumulated history the trend slide draws its composite-score
    series from. This grows on its own as the scheduler actually runs
    week over week against an updated project plan; nothing here is
    backfilled or modeled. Keyed by project_name."""
    by_slug: dict[str, list[tuple[Path, dict]]] = {}
    for path in sorted(OUTPUT_DIR.glob("*_weekly_report_*.json")):
        m = re.match(r"(.+)_weekly_report_\d{8}T\d{6}\.json$", path.name)
        if not m:
            continue
        slug = m.group(1)
        data = json.loads(path.read_text(encoding="utf-8"))
        by_slug.setdefault(slug, []).append((path, data))

    history = {}
    for items in by_slug.values():
        items.sort(key=lambda pair: pair[1]["generated_at"])
        reports_sorted = [d for _, d in items]
        history[reports_sorted[-1]["project_name"]] = reports_sorted
    return history


def _latest_reports() -> list[dict]:
    """Single freshest snapshot per project — what every "current status"
    slide should show. See _report_history() for the full accumulated
    record the trend slide uses instead."""
    return [items[-1] for items in _report_history().values()]


def _month_key(date_str: str) -> str:
    """Which calendar month a report's week belongs to. Not a raw slice
    of report_run_date — a week straddling a month boundary (e.g. Mon
    Jun 29 - Sun Jul 5) is assigned as a whole to whichever month
    contains its Wednesday, via rag_engine.week_calendar_month(), so the
    same week is never split across two different "monthly" groupings."""
    return rag_engine.week_calendar_month(pd.Timestamp(date_str))


def _current_month_reports(history: dict[str, list[dict]]) -> dict[str, list[dict]]:
    """Bounds the trend slide to one real calendar month's worth of
    weekly reports — 4 of them most months, 5 in months with an extra
    scheduled run — rather than every weekly report ever generated. Each
    project is scoped to its own most recent calendar month (by
    report_run_date), so the monthly synthesis always means "this month",
    not an ever-growing accumulation across every month the pipeline has
    run since day one."""
    result = {}
    for project_name, items in history.items():
        if not items:
            continue
        latest_month = _month_key(items[-1]["report_run_date"])
        result[project_name] = [r for r in items if _month_key(r["report_run_date"]) == latest_month]
    return result


def _deterministic_recommendations(reports: list[dict], patterns: list[str]) -> list[str]:
    """No-LLM fallback, built straight from data already computed by
    rag_engine.py / _cross_project_patterns. Used when the model's
    response can't be parsed as a clean JSON array even loosely — this
    happens more with small local models than with a hosted API, and a
    naive line-split of garbled output risks shipping nonsense bullets on
    a VP-facing slide, so a fully deterministic list is the safer floor."""
    recs = []
    for r in reports:
        if r["rag_status"] in ("Red", "Amber") and r.get("red_signals"):
            worst = ", ".join(name.replace("_", " ") for name in r["red_signals"])
            recs.append(
                f"{r['project_name']} ({r['rag_status']}): address {worst} first — these are the "
                "signals currently driving the status."
            )
    for p in patterns[:2]:
        recs.append(f"Portfolio-wide: {p}")
    if not recs:
        recs.append("No project requires urgent action this month; maintain the current monitoring cadence.")
    return recs[:5]


def _generate_recommendations(reports: list[dict], patterns: list[str]) -> list[str]:
    lines = []
    for r in reports:
        lines.append(
            f"- {r['project_name']}: {r['rag_status']} (confidence {r['confidence']}). "
            f"{r['explanation']['narrative']}"
        )
    pattern_block = "\n".join(f"- {p}" for p in patterns) if patterns else "- None detected."
    prompt = (
        "Here is this month's project health portfolio:\n\n" + "\n".join(lines) +
        "\n\nShared patterns already detected across the portfolio (do not repeat these verbatim, "
        "but let them inform your recommendations where relevant):\n" + pattern_block +
        "\n\nWrite 3 to 5 short, specific, action-oriented recommendations for a VP audience, "
        "based only on the data above. At least one recommendation should address a pattern that "
        "spans multiple projects, not just a single project's issue, if the data supports one. "
        "Each recommendation should name the project(s) it applies to. "
        "Output ONLY a fenced ```json code block containing a JSON array of strings, nothing else."
    )
    text = llm_client.generate(
        "You produce crisp, VP-facing portfolio recommendations grounded strictly in the data given.",
        prompt,
    )
    def _valid_list(candidate) -> bool:
        return isinstance(candidate, list) and bool(candidate) and all(isinstance(x, str) for x in candidate)

    fenced = re.search(r"```json\s*(\[.*?\])\s*```", text, re.DOTALL)
    if fenced:
        try:
            parsed = json.loads(fenced.group(1))
            if _valid_list(parsed):
                return parsed[:5]
        except json.JSONDecodeError:
            pass

    # Looser fallback: the model may have produced a JSON array without the
    # fence, or with stray text around it — try the first "[...]" span too.
    bracketed = re.search(r"\[.*\]", text, re.DOTALL)
    if bracketed:
        try:
            parsed = json.loads(bracketed.group(0))
            if _valid_list(parsed):
                return parsed[:5]
        except json.JSONDecodeError:
            pass

    return _deterministic_recommendations(reports, patterns)


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, x, y, w, h, color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.06
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line:
        shape.line.color.rgb = color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _textbox(slide, x, y, w, h, text, size=14, color=SLATE, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return box


def _slide_header(slide, title, subtitle=None):
    _textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.55), title,
             size=32, color=NAVY, bold=True, font="Cambria")
    if subtitle:
        _textbox(slide, Inches(0.6), Inches(1.15), Inches(12), Inches(0.4), subtitle,
                 size=14, color=MUTED)


def _rag_chip(slide, x, y, status, w=Inches(1.5), h=Inches(0.4)):
    color = RAG_COLORS.get(status, RAG_COLORS["Insufficient Data"])
    chip = _rect(slide, x, y, w, h, color)
    tf = chip.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = status
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"
    chip.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return chip


def _period_label(history: dict, reports: list[dict]) -> str:
    """What real date range this synthesis is actually built from — the
    earliest and latest report_run_date across every stored weekly
    report, plus how many weekly runs that spans, per project. Falls
    back to just `reports` if no history dict is supplied."""
    all_reports = [r for items in (history or {}).values() for r in items] or reports
    dates = sorted(r["report_run_date"][:10] for r in all_reports if r.get("report_run_date"))
    if not dates:
        return ""
    max_weeks = max((len(items) for items in (history or {}).values()), default=1)
    week_word = "weekly run" if max_weeks == 1 else "weekly runs"
    if dates[0] == dates[-1]:
        return f"Report period: {dates[0]} ({max_weeks} {week_word} recorded so far)"
    return f"Report period: {dates[0]} to {dates[-1]} (up to {max_weeks} {week_word} recorded so far)"


def add_title_slide(prs, reports, generated_label, history=None):
    slide = _blank_slide(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    _textbox(slide, Inches(0.9), Inches(2.5), Inches(11.5), Inches(0.75),
              "Monthly Portfolio Health Synthesis", size=40, color=WHITE, bold=True, font="Cambria")
    _textbox(slide, Inches(0.9), Inches(3.5), Inches(11.5), Inches(0.6),
              "Cross-project RAG status, synthesized from live project-plan data", size=18, color=ICE)
    _textbox(slide, Inches(0.9), Inches(4.1), Inches(11.5), Inches(0.5),
              f"Generated {generated_label}  |  {len(reports)} project(s)", size=13, color=ICE)
    period = _period_label(history, reports)
    if period:
        _textbox(slide, Inches(0.9), Inches(4.55), Inches(11.5), Inches(0.5),
                  period, size=13, color=ICE)


def _fit_rows(top, bottom, min_row_h, gap, n_items, max_row_h=None):
    """Shared row-layout math for any slide that draws one card/row per
    project. Rather than a fixed row height that silently overflows the
    slide once there are more than ~2-3 projects, this shrinks row height
    to fit whatever's available, and caps the count so rows never drop
    below min_row_h — any remainder is reported via the returned overflow
    count rather than drawn off the bottom of the slide.

    max_row_h caps growth in the other direction: with few projects,
    filling 100% of the available space would stretch each card far
    beyond what its content needs (empty-looking, oversized cards) —
    capping it leaves the extra as plain bottom margin instead."""
    max_rows = max(1, int((bottom - top + gap) / (min_row_h + gap)))
    shown_count = min(n_items, max_rows)
    overflow = n_items - shown_count
    list_bottom = bottom - (Inches(0.35) if overflow else Inches(0))
    row_h = (list_bottom - top + gap) / shown_count - gap
    if max_row_h is not None:
        row_h = min(row_h, max_row_h)
    return shown_count, row_h, list_bottom, overflow


def _add_overflow_note(slide, list_bottom, overflow, left=Inches(0.6), width=Inches(12.1)):
    if overflow:
        _textbox(slide, left, list_bottom + Inches(0.05), width, Inches(0.3),
                  f"+{overflow} more project(s) not shown here — see individual JSON/txt reports.",
                  size=10, color=MUTED)


def add_portfolio_overview_slide(prs, reports):
    slide = _blank_slide(prs)
    _slide_header(slide, "Portfolio Health Overview", "Traffic-light status across all tracked projects")

    left = Inches(0.6)
    width = Inches(12.1)
    gap = Inches(0.2)
    shown_count, row_h, list_bottom, overflow = _fit_rows(
        Inches(1.7), Inches(7.15), Inches(0.85), gap, len(reports), max_row_h=Inches(1.0)
    )

    y = Inches(1.7)
    for r in reports[:shown_count]:
        _rect(slide, left, y, width, row_h, CARD_BG)
        _textbox(slide, left + Inches(0.3), y + Inches(0.1), Inches(6.5), Inches(0.35),
                  r["project_name"], size=15, color=NAVY, bold=True)
        score_txt = f"Composite score: {r['composite_score']}" if r["composite_score"] is not None else "Composite score: n/a"
        _textbox(slide, left + Inches(0.3), y + Inches(0.47), Inches(6.5), Inches(0.32),
                  f"{score_txt}   |   Confidence: {r['confidence']}", size=11, color=MUTED)
        _rag_chip(slide, left + Inches(8.2), y + Inches(0.13), r["rag_status"])
        red_ct = len(r.get("red_signals", []))
        note = f"{red_ct} red signal(s)" if red_ct else "no red signals"
        _textbox(slide, left + Inches(10.0), y + Inches(0.15), Inches(2.0), Inches(0.32),
                  note, size=11, color=MUTED)
        y += row_h + gap

    _add_overflow_note(slide, list_bottom, overflow)


def _cross_project_patterns(reports: list[dict]) -> list[str]:
    """Deterministic (non-LLM) detection of signals that are weak across
    multiple projects at once, so the deck can call out portfolio-wide
    trends rather than just per-project summaries."""
    patterns = []

    weak_by_signal: dict[str, list[str]] = {}
    for r in reports:
        for name, s in r["signals"].items():
            if s and s["score"] < 100:
                weak_by_signal.setdefault(name, []).append(r["project_name"])

    for name, projects in weak_by_signal.items():
        if len(projects) >= 2:
            label = name.replace("_", " ").title()
            patterns.append(f"{label} is a shared weak spot across all {len(projects)} tracked project(s).")

    if reports and not any(r["has_budget_data"] for r in reports):
        patterns.append("No project in this portfolio has budget/cost tracking data in its source plan.")

    sentiment_covered = sum(1 for r in reports if r["signals"].get("stakeholder_sentiment") is not None)
    if 0 < sentiment_covered < len(reports):
        patterns.append(f"Stakeholder sentiment data exists for only {sentiment_covered} of {len(reports)} project(s).")

    return patterns


def _risky_signals(r: dict) -> list[tuple[str, dict]]:
    return [(name, s) for name, s in r["signals"].items() if s and s["score"] < 100]


def add_emerging_risks_slide(prs, reports, patterns: list[str], history: dict = None):
    slide = _blank_slide(prs)
    _slide_header(slide, "Emerging Risks",
                   "Signals currently pulling status down (single snapshot — see Trend Analysis)")

    x = Inches(0.6)
    full_w = Inches(12.1)
    gap = Inches(0.3)

    if patterns:
        _rect(slide, x, Inches(1.65), Inches(12.1), Inches(0.65), ICE)
        _textbox(slide, x + Inches(0.3), Inches(1.75), Inches(11.5), Inches(0.5),
                  "Portfolio pattern: " + "  •  ".join(patterns[:2]), size=12, color=NAVY, bold=True)
        y0 = Inches(2.45)
        bottom = Inches(6.5)
    else:
        y0 = Inches(1.7)
        bottom = Inches(6.75)

    n = len(reports)

    if n <= 3:
        # Card grid: one column per project, full detail on each.
        col_w = (full_w - gap * (n - 1)) / n if n > 0 else full_w
        card_h = bottom - y0
        for i, r in enumerate(reports):
            cx = x + i * (col_w + gap)
            _rect(slide, cx, y0, col_w, card_h, CARD_BG)
            _textbox(slide, cx + Inches(0.3), y0 + Inches(0.25), col_w - Inches(0.6), Inches(0.4),
                      r["project_name"], size=16, color=NAVY, bold=True)
            _rag_chip(slide, cx + Inches(0.3), y0 + Inches(0.7), r["rag_status"], w=Inches(1.3), h=Inches(0.35))

            risky = _risky_signals(r)
            ty = y0 + Inches(1.25)
            if not risky:
                _textbox(slide, cx + Inches(0.3), ty, col_w - Inches(0.6), Inches(0.4),
                          "All signals currently on track.", size=13, color=MUTED)
            for name, s in risky:
                color = RAG_COLORS["Red"] if s["score"] == 20 else RAG_COLORS["Amber"]
                phrase = plain_english_phrase(name, (s["score"], s["raw_metric"]))
                phrase = phrase[0].upper() + phrase[1:]
                _textbox(slide, cx + Inches(0.3), ty, Inches(0.25), Inches(0.35), "●", size=12, color=color)
                _textbox(slide, cx + Inches(0.6), ty, col_w - Inches(1.0), Inches(0.65),
                          phrase, size=12, color=SLATE)
                ty += Inches(0.7)
    else:
        # Too many projects for full detail cards to stay readable —
        # switch to a compact one-row-per-project list, capped so rows
        # never shrink below a legible height.
        shown_count, row_h, list_bottom, overflow = _fit_rows(
            y0, bottom, Inches(0.7), Inches(0.12), len(reports), max_row_h=Inches(1.1)
        )
        y = y0
        for r in reports[:shown_count]:
            _rect(slide, x, y, full_w, row_h, CARD_BG)
            _textbox(slide, x + Inches(0.25), y + Inches(0.1), Inches(5.5), row_h - Inches(0.2),
                      r["project_name"], size=13, color=NAVY, bold=True)
            _rag_chip(slide, x + Inches(6.0), y + Inches(0.08), r["rag_status"], w=Inches(1.1), h=Inches(0.32))
            risky = _risky_signals(r)
            phrases = [plain_english_phrase(name, (s["score"], s["raw_metric"])) for name, s in risky[:2]]
            summary = "; ".join(phrases) or "All signals currently on track."
            summary = summary[0].upper() + summary[1:]
            _textbox(slide, x + Inches(7.3), y + Inches(0.1), Inches(4.5), row_h - Inches(0.2),
                      summary, size=10.5, color=SLATE)
            y += row_h + Inches(0.12)

        _add_overflow_note(slide, list_bottom, overflow)

    history = history or {}
    weeks_recorded = max((len(history.get(r["project_name"], [r])) for r in reports), default=1)
    note_y = Inches(6.65) if patterns else Inches(6.9)
    if weeks_recorded < 2:
        note = ("Note: only one weekly data snapshot exists per project so far; see Trend Analysis for "
                "the real composite-score trend, which populates automatically as more weekly runs accumulate.")
    else:
        note = (f"Note: {weeks_recorded} weekly snapshots recorded for at least one project — see Trend "
                "Analysis for the real composite-score trend across them.")
    _textbox(slide, Inches(0.6), note_y, Inches(12), Inches(0.4), note, size=11, color=MUTED)


def add_milestone_tracker_slide(prs, reports):
    slide = _blank_slide(prs)
    _slide_header(slide, "Milestone Tracker", "Delayed or missed milestones by project")

    gap = Inches(0.2)
    shown_count, row_h, list_bottom, overflow = _fit_rows(
        Inches(1.7), Inches(7.15), Inches(1.3), gap, len(reports), max_row_h=Inches(2.2)
    )

    y = Inches(1.7)
    for r in reports[:shown_count]:
        m = r["signals"].get("milestone_health")
        _rect(slide, Inches(0.6), y, Inches(12.1), row_h, CARD_BG)
        _textbox(slide, Inches(0.9), y + Inches(0.15), Inches(6), Inches(0.35),
                  r["project_name"], size=15, color=NAVY, bold=True)
        if m is None:
            _textbox(slide, Inches(0.9), y + Inches(0.55), Inches(11.5), Inches(0.35),
                      "No milestone data available for this project.", size=12, color=MUTED)
        else:
            chip_status = "Green" if m["score"] == 100 else ("Amber" if m["score"] == 60 else "Red")
            _rag_chip(slide, Inches(9.5), y + Inches(0.15), chip_status, w=Inches(1.2), h=Inches(0.32))
            raw = m["raw_metric"]
            summary, _, names = raw.partition(":")
            _textbox(slide, Inches(0.9), y + Inches(0.55), Inches(11.5), Inches(0.35),
                      summary.strip(), size=13, color=SLATE, bold=True)
            names = names.strip()
            if len(names) > 150:
                names = names[:147].rsplit(",", 1)[0] + "..."
            if names:
                _textbox(slide, Inches(0.9), y + Inches(0.95), Inches(11.5), row_h - Inches(1.0),
                          names, size=10.5, color=MUTED)
        y += row_h + gap

    _add_overflow_note(slide, list_bottom, overflow)


def add_budget_summary_slide(prs, reports):
    slide = _blank_slide(prs)
    _slide_header(slide, "Budget Summary")

    _rect(slide, Inches(0.6), Inches(1.7), Inches(12.1), Inches(1.2), RGBColor(0xFC, 0xF3, 0xE3))
    _textbox(slide, Inches(0.9), Inches(1.95), Inches(11.5), Inches(0.7),
              "No cost/budget column exists in either source project plan. Budget burn cannot be "
              "measured this period — figures below are not fabricated.",
              size=14, color=RGBColor(0x8A, 0x5A, 0x00), bold=True)

    gap = Inches(0.25)
    shown_count, card_h, list_bottom, overflow = _fit_rows(
        Inches(3.1), Inches(7.15), Inches(0.9), gap, len(reports), max_row_h=Inches(1.5)
    )

    y = Inches(3.1)
    for r in reports[:shown_count]:
        b = r["signals"].get("budget_burn")
        _rect(slide, Inches(0.6), y, Inches(12.1), card_h, CARD_BG)
        _textbox(slide, Inches(0.9), y + Inches(0.15), Inches(6), Inches(0.35),
                  r["project_name"], size=15, color=NAVY, bold=True)
        _rag_chip(slide, Inches(9.5), y + Inches(0.15), "N/A", w=Inches(1.2), h=Inches(0.35))
        _textbox(slide, Inches(0.9), y + Inches(0.55), Inches(11.5), card_h - Inches(0.6),
                  (b["raw_metric"] if b else "Signal dropped — no budget data in source; weight redistributed to the other four signals"),
                  size=11, color=MUTED)
        y += card_h + gap

    _add_overflow_note(slide, list_bottom, overflow)


_TREND_STATUS_WORD = {100: "Green", 60: "Amber", 20: "Red"}


def _format_report_date(generated_at: str) -> str:
    """generated_at is stored as 'YYYYMMDDTHHMMSS' (see main.py); render
    it as a short human date label for the trend chart's axis ticks."""
    try:
        return datetime.strptime(generated_at, "%Y%m%dT%H%M%S").strftime("%b %d")
    except ValueError:
        return generated_at


def add_trend_analysis_slide(prs, reports=None, history=None):
    """Two real trends, nothing modeled: the composite RAG score across
    however many weekly runs have actually accumulated per project (real
    once 2+ exist — see _report_history()), and stakeholder sentiment
    over time (real from day one, since comment timestamps are genuine
    dated facts). Every other signal still only has a single current-
    state value per run, so no other composite reconstruction is
    attempted."""
    slide = _blank_slide(prs)
    _slide_header(slide, "Trend Analysis")

    reports = reports or []
    history = history or {}

    _textbox(slide, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.35),
              "Composite RAG Score Over Time (real weekly runs, not modeled)",
              size=16, color=NAVY, bold=True)

    gap = Inches(0.12)
    shown_count, row_h, list_bottom, overflow = _fit_rows(
        Inches(1.95), Inches(3.75), Inches(0.55), gap, len(reports), max_row_h=Inches(0.85)
    )

    y = Inches(1.95)
    for r in reports[:shown_count]:
        proj_hist = history.get(r["project_name"], [r])
        _rect(slide, Inches(0.6), y, Inches(12.1), row_h, CARD_BG)
        _textbox(slide, Inches(0.8), y + Inches(0.08), Inches(2.4), row_h - Inches(0.16),
                  r["project_name"], size=12, color=NAVY, bold=True)
        if len(proj_hist) < 2:
            _textbox(slide, Inches(3.3), y + Inches(0.08), Inches(8.9), row_h - Inches(0.16),
                      f"Only 1 weekly run recorded so far ({_format_report_date(proj_hist[0]['generated_at'])}) — "
                      "a trend needs at least 2 to appear. This fills in automatically as the scheduler keeps running.",
                      size=11, color=MUTED)
        else:
            col_w = Inches(9.0) / len(proj_hist)
            for i, point in enumerate(proj_hist):
                cx = Inches(3.3) + i * col_w
                _rag_chip(slide, cx, y + Inches(0.05), point["rag_status"], w=col_w - Inches(0.15), h=Inches(0.3))
                score_label = "n/a" if point.get("composite_score") is None else f"{point['composite_score']:.0f}"
                _textbox(slide, cx, y + Inches(0.4), col_w - Inches(0.15), Inches(0.25),
                          f"{_format_report_date(point['generated_at'])} · {score_label}", size=8, color=MUTED)
        y += row_h + gap

    _add_overflow_note(slide, list_bottom, overflow)

    trending = [r for r in reports if r.get("sentiment_trend")]
    if not trending:
        _textbox(slide, Inches(0.6), Inches(4.0), Inches(12.1), Inches(0.4),
                  "No project has enough real comment history yet to show a sentiment trend.",
                  size=13, color=MUTED)
        return

    _textbox(slide, Inches(0.6), Inches(4.0), Inches(12.1), Inches(0.4),
              "What's ALSO real: Stakeholder Sentiment Over Time", size=16, color=NAVY, bold=True)
    _textbox(slide, Inches(0.6), Inches(4.4), Inches(12.1), Inches(0.35),
              "Comment timestamps are real, dated facts — re-scoring sentiment using only comments that "
              "existed by each past week is measuring history, not modeling it.",
              size=11, color=MUTED)

    col_labels = ["3 weeks ago", "2 weeks ago", "1 week ago", "This week"]
    sent_gap = Inches(0.15)
    sent_shown, sent_row_h, sent_bottom, sent_overflow = _fit_rows(
        Inches(4.9), Inches(7.15), Inches(0.7), sent_gap, len(trending), max_row_h=Inches(0.85)
    )
    y = Inches(4.9)
    for r in trending[:sent_shown]:
        _rect(slide, Inches(0.6), y, Inches(12.1), sent_row_h, CARD_BG)
        _textbox(slide, Inches(0.8), y + Inches(0.08), Inches(2.6), sent_row_h - Inches(0.16),
                  r["project_name"], size=12, color=NAVY, bold=True)
        col_w = Inches(2.2)
        for i, point in enumerate(r["sentiment_trend"]):
            cx = Inches(3.5) + i * col_w
            word = "N/A" if point["score"] is None else _TREND_STATUS_WORD.get(point["score"], "N/A")
            _rag_chip(slide, cx, y + Inches(0.08), word, w=Inches(1.0), h=Inches(0.3))
            _textbox(slide, cx, y + Inches(0.46), Inches(2.0), Inches(0.3),
                      col_labels[i], size=9, color=MUTED)
        y += sent_row_h + sent_gap

    _add_overflow_note(slide, sent_bottom, sent_overflow)


def add_recommendations_slide(prs, recommendations, subtitle="Generated for a VP audience from this month's portfolio data"):
    slide = _blank_slide(prs)
    _slide_header(slide, "Recommendations", subtitle)

    y = Inches(1.8)
    for i, rec in enumerate(recommendations, start=1):
        number_chip = _rect(slide, Inches(0.6), y, Inches(0.55), Inches(0.55), NAVY)
        tf = number_chip.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(i)
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = WHITE

        _textbox(slide, Inches(1.35), y + Inches(0.02), Inches(11.2), Inches(0.9),
                  rec, size=15, color=SLATE)
        y += Inches(1.0)


def build_deck(reports: list[dict]) -> Path:
    if not reports:
        raise RuntimeError("No JSON reports found in outputs/. Run main.py first.")

    patterns = _cross_project_patterns(reports)
    recommendations = _generate_recommendations(reports, patterns)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    history = _current_month_reports(_report_history())
    generated_label = datetime.now().strftime("%Y-%m-%d %H:%M")
    add_title_slide(prs, reports, generated_label, history)
    add_portfolio_overview_slide(prs, reports)
    add_emerging_risks_slide(prs, reports, patterns, history)
    add_milestone_tracker_slide(prs, reports)
    add_budget_summary_slide(prs, reports)
    add_trend_analysis_slide(prs, reports, history)
    add_recommendations_slide(prs, recommendations)

    OUTPUT_DIR.mkdir(exist_ok=True)
    prs.save(DECK_PATH)
    return DECK_PATH


def main():
    reports = _latest_reports()
    path = build_deck(reports)
    print(f"Deck written to {path}")
    return path


if __name__ == "__main__":
    main()
